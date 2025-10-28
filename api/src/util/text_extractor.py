import logging
from collections.abc import Callable
from dataclasses import dataclass
from functools import cache
from io import BytesIO

import pandas
import pptx
from bs4 import BeautifulSoup
from docx import Document
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text

from src.util.file_util import read_file

logger = logging.getLogger(__name__)


@cache
def get_text_extractor_configs() -> dict:
    html_extractor_config = TextExtractorConfig(extractor=extract_text_from_html)
    xls_extractor_config = TextExtractorConfig(extractor=extract_text_from_xls, read_mode="rb")
    return {
        "docx": TextExtractorConfig(extractor=extract_text_from_docx, read_mode="rb"),
        "pdf": TextExtractorConfig(extractor=extract_text_from_pdf, read_mode="rb"),
        "pptx": TextExtractorConfig(extractor=extract_text_from_pptx, read_mode="rb"),
        "txt": TextExtractorConfig(),
        "html": html_extractor_config,
        "htm": html_extractor_config,
        "rtf": TextExtractorConfig(extractor=extract_text_from_rft),
        "xlsm": xls_extractor_config,
        "xlsx": xls_extractor_config,
    }


class BaseTextExtractorError(Exception):
    def __init__(self, message: str) -> None:
        super().__init__(message)
        self.message = message

    def __str__(self) -> str:
        return self.message


class UnsupportedTextExtractorFileType(BaseTextExtractorError):
    pass


class FileTypeMismatchTextExtractorError(BaseTextExtractorError):
    pass


class UnprocessableFileContent(BaseTextExtractorError):
    pass


class TextExtractorUnknownError(BaseTextExtractorError):
    pass


def extract_text_from_file(
    file_path: str,
    file_type: str | None = None,
    raise_on_error: bool = False,
    char_limit: int | None = None,
) -> str | None:
    """Get text content from a file

    If raise_on_error set to True, the following exceptions will be raised:
        UnsupportedTextExtractorFileType
        FileTypeMismatchTextExtractorError
        UnprocessableFileContent
        TextExtractorUnknownError
    """
    try:
        extracted_text = TextExtractor(file_path, file_type=file_type).get_text()
        if char_limit is not None:
            return extracted_text[:char_limit]
        return extracted_text
    except Exception as e:
        if raise_on_error:
            raise e
        return None


def extract_text_from_html(html_string: str) -> str:
    try:
        return "\n".join(BeautifulSoup(html_string, "html.parser").stripped_strings)
    except Exception as e:
        raise UnprocessableFileContent(
            f"UnprocessableFileContent: Error extracting text from html: {e}"
        ) from e


def extract_text_from_pdf(file_data: bytes) -> str:
    try:
        with BytesIO(file_data) as stream:
            reader = PdfReader(stream)
            number_of_pages = len(reader.pages)
            text_data = []
            for page_number in range(number_of_pages):
                page = reader.pages[page_number]
                if text := page.extract_text().strip():
                    text_data.append(text)
            return "\n".join(text_data)
    except Exception as e:
        raise UnprocessableFileContent(
            f"UnprocessableFileContent: Error extracting text from pdf: {e}"
        ) from e


def extract_text_from_pptx(file_data: bytes) -> str:
    try:
        with BytesIO(file_data) as stream:
            presentation = pptx.Presentation(stream)
            text_runs = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if text := run.text.strip():
                                text_runs.append(text)
            return "\n".join(text_runs)
    except Exception as e:
        raise UnprocessableFileContent(
            f"UnprocessableFileContent: Error extraction text from pptx: {e}"
        ) from e


def extract_text_from_rft(rtf_data: str) -> str:
    try:
        return rtf_to_text(rtf_data)
    except Exception as e:
        raise UnprocessableFileContent(
            f"UnprocessableFileContent: Error extracting text from rtf: {e}"
        ) from e


def extract_text_from_docx(file_data: bytes) -> str:
    try:
        with BytesIO(file_data) as stream:
            return "\n".join(paragraph.text for paragraph in Document(stream).paragraphs)
    except Exception as e:
        raise UnprocessableFileContent(
            f"UnprocessableFileContent: Error extracting text from docx: {e}"
        ) from e


def extract_text_from_xls(file_data: bytes) -> str:
    # Get all rows from every sheet in the xlsx file.
    all_data = {}
    try:
        with BytesIO(file_data) as stream:
            excel_file = pandas.ExcelFile(stream)
            for sheet_name in excel_file.sheet_names:
                xls_data = pandas.read_excel(excel_file, sheet_name=sheet_name, dtype=str)
                all_data[sheet_name] = xls_data.astype(str).values.tolist()
            return _xls_dict_to_string(all_data)
    except Exception as e:
        raise UnprocessableFileContent(
            f"UnprocessableFileContent: Error extracting text from xls: {e}"
        ) from e


def _xls_dict_to_string(xls_dict: dict) -> str:
    xls_data = []
    for sheet_name, rows in xls_dict.items():
        xls_data.append(sheet_name)
        for row in rows:
            xls_data.append(",".join(str(i) for i in row))
        xls_data.append("")
    return "\n".join(xls_data)


@dataclass
class TextExtractorConfig:
    extractor: Callable[[bytes], str] | Callable[[str], str] = lambda data: data
    read_mode: str = "r"


class TextExtractor:
    def __init__(
        self,
        file_path: str,
        file_type: str | None = None,
    ) -> None:
        self.file_path = file_path
        self.file_type = file_type.lower() if file_type else self.file_path.split(".")[-1].lower()
        self._validate_file_type()
        self.config = get_text_extractor_configs()[self.file_type]

    def get_text(self) -> str:
        logger.info(
            "text extract start", extra={"file_path": self.file_path, "file_type": self.file_type}
        )
        try:
            extracted_data = self.config.extractor(self._read_file_data()).strip()
            logger.info(
                "text extract success",
                extra={
                    "file_path": self.file_path,
                    "file_type": self.file_type,
                    "characters_extracted": len(extracted_data),
                },
            )
            return extracted_data
        except BaseTextExtractorError as e:
            raise e
        except Exception as e:
            err = f"TextExtractorUnknownError: Unknown error extracting text from file {self.file_path=} - {self.file_type=}: {e}"
            raise TextExtractorUnknownError(err) from e

    def _read_file_data(self) -> bytes | str:
        try:
            return read_file(self.file_path, self.config.read_mode)
        except Exception as e:
            raise UnprocessableFileContent(
                f"UnprocessableFileContent: Could not read {self.file_type} content: {e}"
            ) from e

    def _validate_file_type(self) -> None:
        if not self.file_path.lower().endswith(self.file_type):
            err = f"FileTypeMismatchTextExtractorError: {self.file_path} must end with {self.file_type}"
            raise FileTypeMismatchTextExtractorError(err)
        if self.file_type not in set(get_text_extractor_configs().keys()):
            err = f"UnsupportedTextExtractorFileType: {self.file_type}"
            raise UnsupportedTextExtractorFileType(err)
